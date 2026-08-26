"""Generate single 3MT slide for PostCare AI."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent / "3MT-postcare-slide.pptx"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x64, 0x74, 0x8B)
TEXT = RGBColor(0x1E, 0x29, 0x3B)
GREEN = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
SKY = RGBColor(0xDB, 0xEA, 0xFE)


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_para(tf, text, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT, space_after=4):
    p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(space_after)
    return p


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Left panel background
    left_bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(6.8), Inches(7.5)
    )
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = NAVY
    left_bg.line.fill.background()

    # Right panel subtle bg
    right_bg = slide.shapes.add_shape(
        1, Inches(6.8), Inches(0), Inches(6.533), Inches(7.5)
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = WHITE
    right_bg.line.fill.background()

    # ── LEFT COLUMN ──
    tb = add_textbox(slide, Inches(0.45), Inches(0.35), Inches(6.0), Inches(0.35))
    set_para(tb.text_frame, "3 MINUTE THESIS · CAPSTONE", 9, True, AMBER)

    tb = add_textbox(slide, Inches(0.45), Inches(0.75), Inches(6.0), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf, "PostCare AI", 40, True, WHITE, space_after=8)
    p = tf.add_paragraph()
    p.text = "AI-guided post-operative wound monitoring with human clinicians in the loop"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    # Problem box
    prob = slide.shapes.add_shape(1, Inches(0.45), Inches(2.35), Inches(6.0), Inches(1.55))
    prob.fill.solid()
    prob.fill.fore_color.rgb = RGBColor(0x1E, 0x29, 0x3B)
    prob.line.color.rgb = AMBER
    prob.line.width = Pt(2)
    tb = add_textbox(slide, Inches(0.6), Inches(2.5), Inches(5.7), Inches(1.3))
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf, "Problem:", 11, True, AMBER, space_after=4)
    p = tf.add_paragraph()
    p.text = (
        "After discharge, wound complications are often caught late. "
        "Patients lack daily guidance; clinicians cannot review every photo manually."
    )
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE

    tb = add_textbox(slide, Inches(0.45), Inches(6.85), Inches(6.0), Inches(0.4))
    set_para(tb.text_frame, "Pakistan Kidney & Liver Institute · Liver Transplant", 10, False, SLATE)

    # ── RIGHT COLUMN ──
    tb = add_textbox(slide, Inches(7.05), Inches(0.35), Inches(5.9), Inches(0.3))
    set_para(tb.text_frame, "SOLUTION PIPELINE", 9, True, SLATE)

    steps = [
        ("1", "Patient check-in", "Photo + pain, symptoms, post-op day"),
        ("2", "MedSigLIP (fine-tuned)", "6 labels: healing, infection risk, urgency…"),
        ("3", "Safety rules + AI agents", "Triage → patient guidance → clinician summary"),
        ("4", "RAG recovery chat", "Vector search + patient case data for Q&A"),
        ("5", "Clinician dashboard", "Priority queue · human review · mark complete"),
    ]
    y = 0.72
    for num, title, sub in steps:
        circle = slide.shapes.add_shape(1, Inches(7.05), Inches(y), Inches(0.28), Inches(0.28))
        circle.fill.solid()
        circle.fill.fore_color.rgb = SKY
        circle.line.fill.background()
        tb = add_textbox(slide, Inches(7.06), Inches(y + 0.02), Inches(0.28), Inches(0.28))
        set_para(tb.text_frame, num, 9, True, BLUE, PP_ALIGN.CENTER)

        tb = add_textbox(slide, Inches(7.45), Inches(y - 0.02), Inches(5.4), Inches(0.55))
        tf = tb.text_frame
        tf.word_wrap = True
        set_para(tf, title, 11, True, TEXT, space_after=2)
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(9)
        p.font.color.rgb = SLATE
        y += 0.62

    # Two boxes
    for i, (heading, items) in enumerate(
        [
            ("MLOps", ["Run 1 → Run 2 lineage", "MLflow + Youden thresholds", "Export to production API"]),
            ("Innovation", ["Vision + language agents", "Decision support, not diagnosis", "Clinician in control"]),
        ]
    ):
        x = 7.05 + i * 2.95
        box = slide.shapes.add_shape(1, Inches(x), Inches(4.05), Inches(2.75), Inches(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        tb = add_textbox(slide, Inches(x + 0.12), Inches(4.12), Inches(2.5), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        set_para(tf, heading, 8, True, SLATE, space_after=4)
        for item in items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(9)
            p.font.color.rgb = TEXT
            p.space_after = Pt(2)

    # Impact
    impact = slide.shapes.add_shape(1, Inches(7.05), Inches(5.55), Inches(5.85), Inches(0.75))
    impact.fill.solid()
    impact.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    impact.line.color.rgb = RGBColor(0xA7, 0xF3, 0xD0)
    tb = add_textbox(slide, Inches(7.2), Inches(5.62), Inches(5.55), Inches(0.65))
    tf = tb.text_frame
    tf.word_wrap = True
    set_para(tf, "Impact: ", 11, True, GREEN, space_after=0)
    tf.paragraphs[0].runs[0].font.bold = True
    p = tf.paragraphs[0]
    p.text = ""
    run_b = p.add_run()
    run_b.text = "Impact: "
    run_b.font.bold = True
    run_b.font.size = Pt(11)
    run_b.font.color.rgb = GREEN
    run_t = p.add_run()
    run_t.text = (
        "Earlier signal detection, scalable remote monitoring, "
        "and structured handoff so surgeons review the right patients first."
    )
    run_t.font.size = Pt(10)
    run_t.font.color.rgb = TEXT

    # Footer
    tb = add_textbox(slide, Inches(7.05), Inches(6.55), Inches(4.5), Inches(0.35))
    set_para(tb.text_frame, "Syed Ashhad Ibrar · PostCare AI", 10, True, TEXT)
    tb = add_textbox(slide, Inches(10.2), Inches(6.55), Inches(2.7), Inches(0.35))
    set_para(tb.text_frame, "Prototype — not clinically validated", 9, False, SLATE, PP_ALIGN.RIGHT)

    prs.save(OUT)
    print(f"Saved: {OUT.resolve()}")


if __name__ == "__main__":
    main()
